from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from parsers.pptx_parser import PptxParser


FIXTURES = Path(__file__).resolve().parents[1] / "fixtures"


def test_pptx_parse_fixture() -> None:
    parser = PptxParser()
    doc = parser.parse(str(FIXTURES / "sample_a.pptx"))
    assert len(doc.segments) >= 2
    assert any(seg.id.startswith("slide1_") for seg in doc.segments)
    assert any(seg.id.startswith("slide2_") for seg in doc.segments)


def test_pptx_can_handle() -> None:
    parser = PptxParser()
    assert parser.can_handle("file.pptx") is True
    assert parser.can_handle("file.docx") is False


def test_pptx_table_extraction(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    rows, cols = 2, 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value 1"
    table.cell(1, 1).text = "Value 2"

    filepath = tmp_path / "table_test.pptx"
    prs.save(str(filepath))

    parser = PptxParser()
    doc = parser.parse(str(filepath))
    targets = [seg.target for seg in doc.segments]
    assert "Header A" in targets
    assert "Header B" in targets
    assert "Value 1" in targets
    assert "Value 2" in targets
    assert any("_tbl_" in seg.id for seg in doc.segments)


def test_pptx_notes_extraction(tmp_path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = "Speaker note line 1"
    notes_frame.add_paragraph().text = "Speaker note line 2"

    filepath = tmp_path / "notes_test.pptx"
    prs.save(str(filepath))

    parser = PptxParser()
    doc = parser.parse(str(filepath))
    targets = [seg.target for seg in doc.segments]
    assert "Speaker note line 1" in targets
    assert "Speaker note line 2" in targets
    assert any("_notes_" in seg.id for seg in doc.segments)
