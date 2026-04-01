from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

from core.models import ParsedDocument, ParseError, Segment, SegmentContext
from parsers.base import BaseParser

_SANITIZE_RE = re.compile(r"[^\w]", re.UNICODE)
_COLLAPSE_RE = re.compile(r"_+")


def _sanitize_shape_name(name: str) -> str:
    result = _SANITIZE_RE.sub("_", name)
    result = _COLLAPSE_RE.sub("_", result).strip("_")
    return result or "unnamed"


class PptxParser(BaseParser):
    name = "PPTX Parser"
    supported_extensions = [".pptx"]
    format_description = "PowerPoint Presentation"

    def can_handle(self, filepath: str) -> bool:
        return Path(filepath).suffix.lower() in self.supported_extensions

    @staticmethod
    def _build_shape_keys(shapes) -> list[str]:
        name_counts: dict[str, int] = {}
        keys: list[str] = []
        for shape in shapes:
            raw_name = getattr(shape, "name", None) or "Shape"
            safe = _sanitize_shape_name(raw_name)
            name_counts[safe] = name_counts.get(safe, 0) + 1
            occ = name_counts[safe]
            keys.append(f"{safe}_{occ}" if occ > 1 else safe)
        return keys

    def parse(self, filepath: str) -> ParsedDocument:
        try:
            presentation = Presentation(filepath)
        except Exception as exc:
            raise ParseError(filepath, str(exc)) from exc

        segments: list[Segment] = []

        def add_segment(segment_id: str, location: str, group: str, text: str) -> None:
            normalized = (text or "").strip()
            if not normalized:
                return
            context = SegmentContext(
                file_path=filepath,
                location=location,
                position=len(segments) + 1,
                group=group,
            )
            segments.append(
                Segment(
                    id=segment_id,
                    source=None,
                    target=normalized,
                    context=context,
                )
            )

        def extract_shape_text(shape, *, slide_index: int, shape_key: str, notes: bool = False) -> None:
            area_name = "Notes" if notes else "Slide"
            group = f"Slide {slide_index} Notes" if notes else f"Slide {slide_index}"
            shape_name = getattr(shape, "name", "Shape")

            if shape.has_text_frame:
                for para_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                    text = (
                        "".join(run.text for run in paragraph.runs)
                        if paragraph.runs
                        else paragraph.text
                    )
                    segment_id = (
                        f"slide{slide_index}_notes_{shape_key}_para{para_index}"
                        if notes
                        else f"slide{slide_index}_{shape_key}_para{para_index}"
                    )
                    location = f"Slide {slide_index} > {area_name} > {shape_name}"
                    add_segment(segment_id, location, group, text)

            if shape.has_table:
                for row_index, row in enumerate(shape.table.rows, start=1):
                    for cell_index, cell in enumerate(row.cells, start=1):
                        segment_id = (
                            f"slide{slide_index}_notes_{shape_key}_tbl_r{row_index}c{cell_index}"
                            if notes
                            else f"slide{slide_index}_{shape_key}_tbl_r{row_index}c{cell_index}"
                        )
                        location = (
                            f"Slide {slide_index} > {area_name} > {shape_name} > "
                            f"Table R{row_index}C{cell_index}"
                        )
                        add_segment(segment_id, location, group, cell.text)

        for slide_index, slide in enumerate(presentation.slides, start=1):
            shape_list = list(slide.shapes)
            shape_keys = self._build_shape_keys(shape_list)
            for shape, shape_key in zip(shape_list, shape_keys):
                extract_shape_text(
                    shape,
                    slide_index=slide_index,
                    shape_key=shape_key,
                )

            if getattr(slide, "has_notes_slide", False):
                try:
                    notes_slide = slide.notes_slide
                except Exception:
                    notes_slide = None
                if notes_slide is not None:
                    notes_shapes = list(notes_slide.shapes)
                    notes_keys = self._build_shape_keys(notes_shapes)
                    for shape, shape_key in zip(notes_shapes, notes_keys):
                        extract_shape_text(
                            shape,
                            slide_index=slide_index,
                            shape_key=shape_key,
                            notes=True,
                        )

        return ParsedDocument(
            segments=segments,
            format_name=self.format_description,
            file_path=filepath,
            metadata={},
            encoding=None,
        )

    def validate(self, filepath: str) -> list[str]:
        errors: list[str] = []
        try:
            _ = Presentation(filepath)
        except Exception as exc:
            errors.append(str(exc))
        return errors
